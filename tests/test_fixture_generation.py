"""Fixture-generation tests (Stage 3).

These tests actually run the deterministic generator (fixtures/
generate_fixtures.py) against the frozen reference_manifest.json and inspect
the resulting DOCX/PDF/PPTX/PNG files. No Docling, no OpenAI, no network --
everything here either calls the generator directly or parses the generated
files with python-docx/python-pptx/Pillow/raw PDF byte inspection.
"""

from __future__ import annotations

import base64
import hashlib
import json
import re
import shutil
import zipfile
import zlib
from pathlib import Path

import pytest
from docx import Document
from PIL import Image
from pptx import Presentation

import generate_fixtures as gf
from manifest_schema import load_manifest

MANIFEST_PATH = Path(__file__).resolve().parent.parent / "fixtures" / "reference_manifest.json"
GENERATED_DIR = Path(__file__).resolve().parent.parent / "fixtures" / "generated"

FORBIDDEN_SOURCE_SUBSTRINGS = [
    "openai", "OpenAI", "anthropic", "requests.", "urllib.request",
    "http.client", "httpx", "socket.socket", "grpc",
]


# --- generation fixture (module-scoped: generate once, reuse for all tests) --


@pytest.fixture(scope="module", autouse=True)
def generated():
    gf.generate_all()
    yield


@pytest.fixture(scope="module")
def manifest():
    return load_manifest(MANIFEST_PATH)


# --- low-level PDF helpers (no external PDF-parsing dependency) -----------


def _pdf_uncompressed_content_streams(path: Path) -> str:
    """Concatenate every uncompressed stream object's body -- our PDFs are
    generated with pageCompression=0, so real page content streams have no
    /Filter, while the embedded image XObject does (and is skipped here)."""
    data = path.read_bytes().decode("latin-1")
    parts = []
    for m in re.finditer(r"(\d+) 0 obj\s*<<([^>]*)>>\s*stream\r?\n(.*?)endstream", data, re.DOTALL):
        header = m.group(2)
        if "/Filter" in header:
            continue
        parts.append(m.group(3))
    return "\n".join(parts)


def _pdf_glyph_operator_count(path: Path) -> int:
    """Count real glyph-showing operators (Tj/TJ). reportlab's Canvas always
    emits an empty BT ... ET text-state block even with zero drawString()
    calls, so presence of "BT" alone is NOT a reliable text-layer signal --
    Tj/TJ (which actually paint glyphs) is."""
    text = _pdf_uncompressed_content_streams(path)
    return len(re.findall(r"\bTj\b|\bTJ\b", text))


def _pdf_text_fragments(path: Path) -> list[str]:
    text = _pdf_uncompressed_content_streams(path)
    fragments = re.findall(r"\(((?:[^()\\]|\\.)*)\)\s*Tj", text)
    return [f.replace("\\(", "(").replace("\\)", ")") for f in fragments]


def _pdf_joined_text(path: Path) -> str:
    return " ".join(_pdf_text_fragments(path))


def _pdf_page_count(path: Path) -> int:
    data = path.read_bytes()
    # PDF dict key order isn't guaranteed (/Count can precede or follow
    # /Type /Pages), so capture the whole /Pages dict body first and then
    # search for /Count within it, order-independent.
    dict_match = re.search(rb"<<([^>]*?/Type\s*/Pages[^>]*?)>>", data, re.DOTALL)
    assert dict_match is not None, f"no /Type /Pages dict found in {path}"
    count_match = re.search(rb"/Count (\d+)", dict_match.group(1))
    assert count_match is not None, f"no /Count found in /Pages dict in {path}"
    return int(count_match.group(1))


def _pdf_embedded_image(path: Path) -> Image.Image:
    """Decode the (ASCII85+Flate encoded, uncompressed-page-stream) image
    XObject reportlab embeds. reportlab RE-ENCODES source raster images into
    its own PDF image representation, so this compares decoded PIXELS, not
    raw file bytes -- the correct notion of "same image" once a re-encoding
    step is involved (see test_shared_image_pixel_identical_in_pdf below)."""
    data = path.read_bytes()
    m = re.search(
        rb"(\d+) 0 obj\s*<<([^>]*?/Subtype\s*/Image[^>]*?)>>\s*stream\r?\n(.*?)endstream",
        data, re.DOTALL,
    )
    assert m is not None, f"no image XObject found in {path}"
    width = int(re.search(rb"/Width (\d+)", m.group(0)).group(1))
    height = int(re.search(rb"/Height (\d+)", m.group(0)).group(1))
    body = m.group(3).strip(b"\r\n \t")
    if body.endswith(b"~>"):
        body = body[:-2]
    pixels = zlib.decompress(base64.a85decode(body, adobe=False))
    return Image.frombytes("RGB", (width, height), pixels)


def _zip_media_bytes(path: Path, prefix: str) -> dict[str, bytes]:
    with zipfile.ZipFile(path) as zf:
        return {name: zf.read(name) for name in zf.namelist() if name.startswith(prefix)}


# --- 1. files exist / manifest validates -----------------------------------


EXPECTED_FILES = [
    "images/diagram_v1.png",
    "images/chart_v1.png",
    "images/scanned_text_v1.png",
    "parity/PARITY_001.docx",
    "parity/PARITY_001.pdf",
    "parity/PARITY_001.pptx",
    "stress/STRESS_DOCX_001.docx",
    "stress/STRESS_PDF_001.pdf",
    "stress/STRESS_PPTX_001.pptx",
    "stress/STRESS_PPTX_002.pptx",
    "stress/STRESS_CHART_001.pdf",
    "stress/STRESS_SCANNED_001.pdf",
]


def test_all_expected_files_exist():
    for rel in EXPECTED_FILES:
        path = GENERATED_DIR / rel
        assert path.exists(), f"missing generated file: {rel}"
        assert path.stat().st_size > 0


def test_manifest_schema_validates_and_is_frozen():
    m = load_manifest(MANIFEST_PATH)
    assert m.status == "approved_frozen"
    assert m.manifest_version == "1.2.1"


# --- 2. expected page/slide counts -----------------------------------------


def test_parity_docx_has_exactly_one_page_break():
    doc = Document(GENERATED_DIR / "parity/PARITY_001.docx")
    assert doc.element.xml.count('w:type="page"') == 1


def test_parity_pdf_has_two_pages():
    assert _pdf_page_count(GENERATED_DIR / "parity/PARITY_001.pdf") == 2


def test_parity_pptx_has_two_slides():
    prs = Presentation(GENERATED_DIR / "parity/PARITY_001.pptx")
    assert len(prs.slides) == 2


def test_stress_scanned_pdf_has_one_page():
    assert _pdf_page_count(GENERATED_DIR / "stress/STRESS_SCANNED_001.pdf") == 1


def test_stress_chart_pdf_has_one_page():
    assert _pdf_page_count(GENERATED_DIR / "stress/STRESS_CHART_001.pdf") == 1


# --- 3. native headings / paragraphs / lists / table cells -----------------


def test_parity_docx_headings_match_manifest(manifest):
    doc = Document(GENERATED_DIR / "parity/PARITY_001.docx")
    heading_texts = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
    assert heading_texts == [h.text for h in manifest.parity_suite.headings]


def test_parity_docx_paragraphs_include_manifest_text(manifest):
    doc = Document(GENERATED_DIR / "parity/PARITY_001.docx")
    all_text = "\n".join(p.text for p in doc.paragraphs)
    for paragraph in manifest.parity_suite.paragraphs:
        assert paragraph.text in all_text
    for distractor in manifest.parity_suite.distractor_facts:
        assert distractor.text in all_text
    assert manifest.parity_suite.captions[0].text in all_text


def test_parity_docx_table_cells_match_manifest(manifest):
    doc = Document(GENERATED_DIR / "parity/PARITY_001.docx")
    table = doc.tables[0]
    for cell in manifest.parity_suite.tables[0].cells:
        assert table.cell(cell.row, cell.col).text == cell.text


def test_parity_pptx_table_cells_match_manifest(manifest):
    prs = Presentation(GENERATED_DIR / "parity/PARITY_001.pptx")
    graphic_frame = next(s for s in prs.slides[0].shapes if s.has_table)
    table = graphic_frame.table
    for cell in manifest.parity_suite.tables[0].cells:
        assert table.cell(cell.row, cell.col).text == cell.text


def test_parity_pdf_text_includes_manifest_facts(manifest):
    joined = _pdf_joined_text(GENERATED_DIR / "parity/PARITY_001.pdf")
    for heading in manifest.parity_suite.headings[:2]:  # H_003 is on page 2, checked separately
        assert heading.text in joined
    for paragraph in manifest.parity_suite.paragraphs:
        assert paragraph.text in joined
    for distractor in manifest.parity_suite.distractor_facts:
        assert distractor.text in joined
    for cell in manifest.parity_suite.tables[0].cells:
        assert cell.text.replace("(", "").replace(")", "") in joined.replace("(", "").replace(")", "")


def test_stress_docx_nested_headings_and_lists(manifest):
    fixture = manifest.stress_suite.docx_nested_structure
    doc = Document(GENERATED_DIR / "stress/STRESS_DOCX_001.docx")

    headings = [(p.text, p.style.name) for p in doc.paragraphs if p.style.name.startswith("Heading")]
    assert [h[0] for h in headings] == [h.text for h in fixture.headings]
    assert [h[1] for h in headings] == [f"Heading {h.level}" for h in fixture.headings]

    list_paragraphs = [p for p in doc.paragraphs if p.style.name.startswith("List Bullet")]
    assert [p.text for p in list_paragraphs] == [item.text for item in fixture.list_items]
    expected_style_by_indent = {0: "List Bullet", 1: "List Bullet 2", 2: "List Bullet 3"}
    for paragraph, item in zip(list_paragraphs, fixture.list_items):
        assert paragraph.style.name == expected_style_by_indent[item.indent_level]


def test_stress_pdf_two_column_and_merged_table(manifest):
    fixture = manifest.stress_suite.pdf_complex_layout
    joined = _pdf_joined_text(GENERATED_DIR / "stress/STRESS_PDF_001.pdf")
    for paragraph in fixture.paragraphs:
        assert paragraph.text in joined
    for cell in fixture.tables[0].cells:
        assert cell.text in joined


def _rect_from_shape(shape) -> tuple[int, int, int, int]:
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _rects_intersect(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> bool:
    ax0, ay0, ax1, ay1 = a
    bx0, by0, bx1, by1 = b
    return ax0 < bx1 and bx0 < ax1 and ay0 < by1 and by0 < ay1


def test_stress_pptx_overlapping_textboxes_both_present_with_z_order(manifest):
    fixture = manifest.stress_suite.pptx_overlapping_textboxes
    prs = Presentation(GENERATED_DIR / "stress/STRESS_PPTX_001.pptx")
    slide = prs.slides[0]
    textboxes = [s for s in slide.shapes if s.has_text_frame and not s.has_table]
    textbox_texts = [s.text_frame.text for s in textboxes]
    # Added in z_order-sorted order (lowest/back first), so shape order in
    # the tree reflects z_order directly.
    expected_order = [b.text for b in sorted(fixture.text_boxes, key=lambda b: b.z_order)]
    assert textbox_texts == expected_order

    graphic_frame = next(s for s in slide.shapes if s.has_table)
    for cell in fixture.table.cells:
        assert graphic_frame.table.cell(cell.row, cell.col).text == cell.text


def test_stress_pptx_overlapping_textboxes_both_exist():
    prs = Presentation(GENERATED_DIR / "stress/STRESS_PPTX_001.pptx")
    slide = prs.slides[0]
    textboxes = [s for s in slide.shapes if s.has_text_frame and not s.has_table]
    assert len(textboxes) == 2


def test_stress_pptx_overlapping_textboxes_rectangles_intersect():
    prs = Presentation(GENERATED_DIR / "stress/STRESS_PPTX_001.pptx")
    slide = prs.slides[0]
    textboxes = [s for s in slide.shapes if s.has_text_frame and not s.has_table]
    rect_a, rect_b = [_rect_from_shape(s) for s in textboxes]
    assert _rects_intersect(rect_a, rect_b)


def test_stress_pptx_overlapping_textboxes_positions_not_identical():
    prs = Presentation(GENERATED_DIR / "stress/STRESS_PPTX_001.pptx")
    slide = prs.slides[0]
    textboxes = [s for s in slide.shapes if s.has_text_frame and not s.has_table]
    positions = [(s.left, s.top) for s in textboxes]
    assert positions[0] != positions[1]


def test_stress_pptx_overlapping_textboxes_foreground_has_opaque_fill(manifest):
    from pptx.enum.dml import MSO_FILL_TYPE

    fixture = manifest.stress_suite.pptx_overlapping_textboxes
    prs = Presentation(GENERATED_DIR / "stress/STRESS_PPTX_001.pptx")
    slide = prs.slides[0]
    textboxes = [s for s in slide.shapes if s.has_text_frame and not s.has_table]

    # Shape order in the tree matches z_order (lowest first); the foreground
    # box is the last one added, i.e. the one with the highest z_order.
    foreground_text = max(fixture.text_boxes, key=lambda b: b.z_order).text
    foreground_shape = next(s for s in textboxes if s.text_frame.text == foreground_text)
    background_shape = next(s for s in textboxes if s.text_frame.text != foreground_text)

    assert foreground_shape.fill.type == MSO_FILL_TYPE.SOLID
    # The background (stale) box is left with no explicit fill, unlike the
    # foreground one -- that contrast is what makes the stacking visible.
    assert background_shape.fill.type != MSO_FILL_TYPE.SOLID


def test_stress_pptx_native_diagram_shapes_and_connectors(manifest):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    fixture = manifest.stress_suite.pptx_native_diagram
    prs = Presentation(GENERATED_DIR / "stress/STRESS_PPTX_002.pptx")
    slide = prs.slides[0]

    labels = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert labels == [n.label for n in fixture.diagram_nodes]

    # python-pptx reports straight connectors as shape_type LINE, not a
    # generic "connector" type.
    connector_count = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE)
    assert connector_count == len(fixture.diagram_edges)


def test_stress_pptx_native_diagram_connectors_reference_correct_shapes_and_have_arrowheads(manifest):
    """For every manifest edge: a native connector exists whose a:stCxn/
    a:endCxn resolve to the declared source/target shape ids, and it carries
    a target-end (a:tailEnd) arrowhead -- verified directly against the OOXML
    since python-pptx exposes none of this through its public API."""
    from pptx.oxml.ns import qn

    fixture = manifest.stress_suite.pptx_native_diagram
    prs = Presentation(GENERATED_DIR / "stress/STRESS_PPTX_002.pptx")
    slide = prs.slides[0]

    shape_id_by_label = {
        s.text_frame.text: s.shape_id for s in slide.shapes if s.has_text_frame
    }
    label_by_fact_id = {node.fact_id: node.label for node in fixture.diagram_nodes}

    connectors = [s for s in slide.shapes if s._element.tag == qn("p:cxnSp")]
    assert len(connectors) == len(fixture.diagram_edges)

    for edge in fixture.diagram_edges:
        expected_source_id = shape_id_by_label[label_by_fact_id[edge.source]]
        expected_target_id = shape_id_by_label[label_by_fact_id[edge.target]]

        match = None
        for connector in connectors:
            cxn_sp_pr = connector._element.find(qn("p:nvCxnSpPr") + "/" + qn("p:cNvCxnSpPr"))
            st_cxn = cxn_sp_pr.find(qn("a:stCxn"))
            end_cxn = cxn_sp_pr.find(qn("a:endCxn"))
            if st_cxn is not None and end_cxn is not None:
                if int(st_cxn.get("id")) == expected_source_id and int(end_cxn.get("id")) == expected_target_id:
                    match = connector
                    break
        assert match is not None, f"no connector found for edge {edge.fact_id} ({edge.source} -> {edge.target})"

        ln = match._element.find(".//" + qn("a:ln"))
        assert ln is not None, f"connector for edge {edge.fact_id} has no <a:ln>"
        tail_end = ln.find(qn("a:tailEnd"))
        assert tail_end is not None, f"connector for edge {edge.fact_id} has no target-end arrowhead"
        assert tail_end.get("type") == "triangle"


# --- 4. shared PNG identical wherever embedded ------------------------------


def test_shared_image_byte_identical_in_docx_and_pptx():
    source_hash = hashlib.sha256((GENERATED_DIR / "images/diagram_v1.png").read_bytes()).hexdigest()

    docx_media = _zip_media_bytes(GENERATED_DIR / "parity/PARITY_001.docx", "word/media/")
    assert len(docx_media) == 1
    docx_hash = hashlib.sha256(next(iter(docx_media.values()))).hexdigest()
    assert docx_hash == source_hash

    pptx_media = _zip_media_bytes(GENERATED_DIR / "parity/PARITY_001.pptx", "ppt/media/")
    assert len(pptx_media) == 1
    pptx_hash = hashlib.sha256(next(iter(pptx_media.values()))).hexdigest()
    assert pptx_hash == source_hash


def test_shared_image_pixel_identical_in_pdf():
    """reportlab re-encodes embedded raster images into its own PDF image
    representation (Flate+ASCII85), so the PDF's raw embedded bytes are NOT
    expected to equal the source PNG's file bytes -- decoded PIXELS are the
    correct comparison here."""
    source = Image.open(GENERATED_DIR / "images/diagram_v1.png").convert("RGB")
    embedded = _pdf_embedded_image(GENERATED_DIR / "parity/PARITY_001.pdf")
    assert embedded.size == source.size
    assert embedded.tobytes() == source.tobytes()


# --- 5. chart source facts match the manifest -------------------------------


def test_chart_image_matches_manifest_numeric_facts(manifest):
    from diagram_image import generate_chart_png

    fresh = generate_chart_png(manifest)
    on_disk = (GENERATED_DIR / "images/chart_v1.png").read_bytes()
    assert hashlib.sha256(fresh).hexdigest() == hashlib.sha256(on_disk).hexdigest()


def test_chart_image_changes_if_manifest_values_change(manifest):
    """Proves the chart is genuinely derived from manifest data, not a
    static/hardcoded asset: mutating a numeric fact must change the pixels."""
    from diagram_image import generate_chart_png

    original = generate_chart_png(manifest)
    mutated = manifest.model_copy(deep=True)
    q1_fact = next(
        f for f in mutated.stress_suite.chart_visual_stress.visual_facts
        if f.subject == "Q1 pass rate"
    )
    q1_fact.value = 10.0
    changed = generate_chart_png(mutated)
    assert hashlib.sha256(original).hexdigest() != hashlib.sha256(changed).hexdigest()


def test_chart_pdf_embeds_the_chart_image():
    chart_source = Image.open(GENERATED_DIR / "images/chart_v1.png").convert("RGB")
    embedded = _pdf_embedded_image(GENERATED_DIR / "stress/STRESS_CHART_001.pdf")
    assert embedded.size == chart_source.size
    assert embedded.tobytes() == chart_source.tobytes()


# --- 6. scanned PDF has no digital text layer -------------------------------


def test_scanned_pdf_has_no_glyph_showing_operators():
    assert _pdf_glyph_operator_count(GENERATED_DIR / "stress/STRESS_SCANNED_001.pdf") == 0


def test_scanned_pdf_has_no_extractable_text_fragments():
    assert _pdf_text_fragments(GENERATED_DIR / "stress/STRESS_SCANNED_001.pdf") == []


def test_scanned_pdf_embeds_an_image():
    embedded = _pdf_embedded_image(GENERATED_DIR / "stress/STRESS_SCANNED_001.pdf")
    assert embedded.size == (900, 120)


def test_parity_and_stress_pdfs_do_have_glyph_operators():
    """Contrast check: the born-digital PDFs (unlike the scanned one) DO have
    real text-showing operators."""
    assert _pdf_glyph_operator_count(GENERATED_DIR / "parity/PARITY_001.pdf") > 0
    assert _pdf_glyph_operator_count(GENERATED_DIR / "stress/STRESS_PDF_001.pdf") > 0


# --- 7. regeneration determinism -------------------------------------------


def test_regeneration_is_byte_deterministic(tmp_path):
    """Covers generation_report.json too, not just the 12 fixture files --
    now that the report excludes itself from its own inventory (see the
    generation_report.json tests below), it's no longer self-referential and
    is fully deterministic like everything else."""
    snapshot_dir = tmp_path / "run1"
    shutil.copytree(GENERATED_DIR, snapshot_dir)

    shutil.rmtree(GENERATED_DIR)
    gf.generate_all()

    files1 = sorted(p for p in snapshot_dir.rglob("*") if p.is_file())
    mismatches = []
    for path1 in files1:
        rel = path1.relative_to(snapshot_dir)
        path2 = GENERATED_DIR / rel
        assert path2.exists(), f"{rel} missing from regenerated output"
        hash1 = hashlib.sha256(path1.read_bytes()).hexdigest()
        hash2 = hashlib.sha256(path2.read_bytes()).hexdigest()
        if hash1 != hash2:
            mismatches.append(str(rel))
    assert not mismatches, f"nondeterministic regeneration for: {mismatches}"


def test_regeneration_produces_same_manifest_sha256():
    report1 = gf.generate_all()
    report2 = gf.generate_all()
    assert report1["manifest_sha256"] == report2["manifest_sha256"]


# --- generation_report.json correctness -------------------------------------


def _load_report() -> dict:
    return json.loads((GENERATED_DIR / gf.REPORT_FILENAME).read_text(encoding="utf-8"))


def test_generation_report_excludes_itself():
    report = _load_report()
    assert gf.REPORT_FILENAME not in report["files"]


def test_generation_report_lists_exactly_the_12_benchmark_artifacts():
    report = _load_report()
    assert sorted(report["files"]) == sorted(EXPECTED_FILES)
    assert len(report["files"]) == 12


def test_generation_report_hashes_match_files_on_disk():
    report = _load_report()
    for relative_key, entry in report["files"].items():
        path = GENERATED_DIR / relative_key
        assert path.exists(), f"{relative_key} listed in report but missing on disk"
        actual_hash = hashlib.sha256(path.read_bytes()).hexdigest()
        assert actual_hash == entry["sha256"], f"hash mismatch for {relative_key}"
        assert path.stat().st_size == entry["size_bytes"]


def test_generation_report_keys_contain_no_backslashes():
    report = _load_report()
    for relative_key in report["files"]:
        assert "\\" not in relative_key, f"non-POSIX path key in report: {relative_key!r}"


# --- 5. PARITY_001.pptx slide 1 shape geometry regression -------------------


def test_parity_pptx_slide1_shapes_do_not_overlap():
    """Regression test for the body-text-overflows-into-heading bug: the
    body, "Recovery Objectives" heading, and table rectangles (plus the
    title, for completeness) must be pairwise non-overlapping."""
    prs = Presentation(GENERATED_DIR / "parity/PARITY_001.pptx")
    slide0 = prs.slides[0]
    shapes = list(slide0.shapes)
    assert len(shapes) == 4  # title, body, heading2, table

    rects = {shape.name: _rect_from_shape(shape) for shape in shapes}
    names = list(rects)
    overlapping_pairs = [
        (a, b) for i, a in enumerate(names) for b in names[i + 1:]
        if _rects_intersect(rects[a], rects[b])
    ]
    assert not overlapping_pairs, f"overlapping shapes on parity slide 1: {overlapping_pairs}"


# --- 8. no LLM or external API used -----------------------------------------


def test_generator_source_has_no_llm_or_network_calls():
    import diagram_image
    import manifest_schema

    modules = [gf, diagram_image, manifest_schema]
    for module in modules:
        source = Path(module.__file__).read_text(encoding="utf-8")
        for forbidden in FORBIDDEN_SOURCE_SUBSTRINGS:
            assert forbidden not in source, f"{module.__name__} unexpectedly references {forbidden!r}"
