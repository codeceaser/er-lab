"""Deterministic fixture generator.

Reads the frozen reference_manifest.json and produces the parity suite
(DOCX/PDF/PPTX, same semantic content, same shared image) and the stress
suite (format-specific edge cases) under fixtures/generated/. No LLM calls,
no network access, no randomness -- every fixture is built from fixed
manifest content and fixed layout/coordinate choices.

Run: python fixtures/generate_fixtures.py
(from the project root, or from inside fixtures/ -- both work, since this
module only uses relative imports of its sibling modules via sys.path
insertion, matching this project's existing script convention.)
"""

from __future__ import annotations

import io
import json
import sys
import zipfile
from datetime import datetime, timezone
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from diagram_image import generate_chart_png, generate_diagram_png, generate_scanned_text_png
from manifest_schema import ReferenceManifest, load_manifest, load_manifest_raw
from rich.console import Console

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))
from ingestion_bench.canonical import compute_manifest_sha256  # noqa: E402

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches as PptxInches, Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen.canvas import Canvas
from reportlab.platypus import Image as RLImage, PageBreak, Paragraph as RLParagraph, SimpleDocTemplate, Spacer, Table, TableStyle

console = Console()

MANIFEST_PATH = Path(__file__).resolve().parent / "reference_manifest.json"
GENERATED_DIR = Path(__file__).resolve().parent / "generated"
PARITY_DIR = GENERATED_DIR / "parity"
STRESS_DIR = GENERATED_DIR / "stress"
IMAGES_DIR = GENERATED_DIR / "images"
REPORT_FILENAME = "generation_report.json"

# Fixed timestamps used everywhere a library would otherwise default to
# "now" -- the single source of the "no wall-clock in any fixture" rule.
FIXED_DOC_DATETIME = datetime(2026, 1, 1, 0, 0, 0, tzinfo=timezone.utc)
FIXED_ZIP_DATETIME = (2026, 1, 1, 0, 0, 0)


# --- determinism helpers ---------------------------------------------------


def _invariant_canvasmaker(*args, **kwargs):
    kwargs["invariant"] = 1
    return Canvas(*args, **kwargs)


def normalize_zip_timestamps(path: Path) -> None:
    """DOCX/PPTX are ZIP packages; zipfile stamps each entry with the current
    local time unless told otherwise. Rewrite every entry with a fixed
    date_time (content, order, and per-entry compression unchanged) so the
    resulting file hashes identically across reruns regardless of wall-clock
    time."""
    with zipfile.ZipFile(path, "r") as zf:
        entries = [(info, zf.read(info.filename)) for info in zf.infolist()]

    with zipfile.ZipFile(path, "w") as zf:
        for info, data in entries:
            new_info = zipfile.ZipInfo(info.filename, date_time=FIXED_ZIP_DATETIME)
            new_info.compress_type = info.compress_type
            new_info.external_attr = info.external_attr
            new_info.create_system = info.create_system
            zf.writestr(new_info, data)


def set_docx_core_properties(document: Document) -> None:
    props = document.core_properties
    props.author = "er-lab-ingestion-bench"
    props.last_modified_by = "er-lab-ingestion-bench"
    props.created = FIXED_DOC_DATETIME
    props.modified = FIXED_DOC_DATETIME
    props.title = ""
    props.subject = ""
    props.comments = ""
    props.category = ""
    props.revision = 1


def set_pptx_core_properties(presentation: Presentation) -> None:
    props = presentation.core_properties
    props.author = "er-lab-ingestion-bench"
    props.last_modified_by = "er-lab-ingestion-bench"
    props.created = FIXED_DOC_DATETIME
    props.modified = FIXED_DOC_DATETIME
    props.title = ""
    props.subject = ""
    props.comments = ""
    props.category = ""
    props.revision = 1


def save_docx(document: Document, path: Path) -> None:
    set_docx_core_properties(document)
    document.save(str(path))
    normalize_zip_timestamps(path)


def save_pptx(presentation: Presentation, path: Path) -> None:
    set_pptx_core_properties(presentation)
    presentation.save(str(path))
    normalize_zip_timestamps(path)


def add_target_arrowhead(connector, arrow_type: str = "triangle") -> None:
    """python-pptx has no public API for line-end arrowheads, so this patches
    the OOXML directly: adds <a:tailEnd type="..."/> to the connector's
    <a:ln> line-properties element (the "tail" end is the line's target/end
    point -- the "head" end is its source/start point). Inserted before any
    existing <a:extLst> to respect CT_LineProperties' element ordering."""
    ln = connector.line._get_or_add_ln()
    tail_end = ln.makeelement(qn("a:tailEnd"), {"type": arrow_type})
    ext_lst = ln.find(qn("a:extLst"))
    if ext_lst is not None:
        ext_lst.addprevious(tail_end)
    else:
        ln.append(tail_end)


# --- parity suite ------------------------------------------------------


def _parity_unit0_paragraphs(manifest: ReferenceManifest) -> list[str]:
    """Body paragraphs + distractor facts for unit 0, in their declared
    order_hint order -- the manifest is the only source of this ordering."""
    entries: list[tuple[int, str]] = []
    for paragraph in manifest.parity_suite.paragraphs:
        entries.append((paragraph.expected_location.order_hint, paragraph.text))
    for distractor in manifest.parity_suite.distractor_facts:
        entries.append((distractor.expected_location.order_hint, distractor.text))
    entries.sort(key=lambda item: item[0])
    return [text for _, text in entries]


def generate_parity_docx(manifest: ReferenceManifest, image_bytes: bytes, out_path: Path) -> None:
    suite = manifest.parity_suite
    doc = Document()

    # Unit 0
    doc.add_heading(suite.headings[0].text, level=1)  # H_001 title
    for text in _parity_unit0_paragraphs(manifest):
        doc.add_paragraph(text)
    doc.add_heading(suite.headings[1].text, level=2)  # H_002
    table_fact = suite.tables[0]
    table = doc.add_table(rows=table_fact.n_rows, cols=table_fact.n_cols)
    table.style = "Table Grid"
    for cell in table_fact.cells:
        target = table.cell(cell.row, cell.col)
        target.text = cell.text
        if cell.is_header:
            for run in target.paragraphs[0].runs:
                run.bold = True

    # Hard unit boundary
    doc.add_page_break()

    # Unit 1
    doc.add_heading(suite.headings[2].text, level=2)  # H_003
    picture_paragraph = doc.add_paragraph()
    picture_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    picture_paragraph.add_run().add_picture(io.BytesIO(image_bytes), width=Inches(6.0))
    caption_paragraph = doc.add_paragraph(suite.captions[0].text)
    caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

    save_docx(doc, out_path)


def generate_parity_pdf(manifest: ReferenceManifest, image_bytes: bytes, out_path: Path) -> None:
    suite = manifest.parity_suite
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("H1Fixed", parent=styles["Heading1"])
    h2 = ParagraphStyle("H2Fixed", parent=styles["Heading2"])
    body = styles["BodyText"]

    story = []
    story.append(RLParagraph(suite.headings[0].text, h1))
    for text in _parity_unit0_paragraphs(manifest):
        story.append(RLParagraph(text, body))
    story.append(Spacer(1, 12))
    story.append(RLParagraph(suite.headings[1].text, h2))

    table_fact = suite.tables[0]
    grid: list[list[str]] = [["" for _ in range(table_fact.n_cols)] for _ in range(table_fact.n_rows)]
    header_cells: set[tuple[int, int]] = set()
    for cell in table_fact.cells:
        grid[cell.row][cell.col] = cell.text
        if cell.is_header:
            header_cells.add((cell.row, cell.col))
    pdf_table = Table(grid)
    style_commands = [
        ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]
    for row, col in header_cells:
        style_commands.append(("FONTNAME", (col, row), (col, row), "Helvetica-Bold"))
    pdf_table.setStyle(TableStyle(style_commands))
    story.append(pdf_table)

    story.append(PageBreak())
    story.append(RLParagraph(suite.headings[2].text, h2))
    story.append(Spacer(1, 6))
    story.append(RLImage(io.BytesIO(image_bytes), width=6.0 * inch, height=6.0 * inch * (160 / 700)))
    story.append(Spacer(1, 6))
    story.append(RLParagraph(suite.captions[0].text, body))

    doc = SimpleDocTemplate(
        str(out_path), pagesize=LETTER, pageCompression=0,
        title="", author="", subject="",
    )
    doc.build(story, canvasmaker=_invariant_canvasmaker)


def generate_parity_pptx(manifest: ReferenceManifest, image_bytes: bytes, out_path: Path) -> None:
    suite = manifest.parity_suite
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Slide 0 (unit 0) -- fixed vertical budget (slide height 7.5in) so
    # body/heading2/table never overlap regardless of how many lines the
    # (compact, 11pt) body paragraphs wrap to:
    #   title:     0.30 - 0.90
    #   body:      1.00 - 3.30  (2.3in, comfortably fits 7 paragraphs @ 11pt)
    #   heading2:  3.50 - 3.95
    #   table:     4.15 - 6.15
    slide0 = prs.slides.add_slide(blank_layout)
    title_box = slide0.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.6))
    title_tf = title_box.text_frame
    title_tf.text = suite.headings[0].text
    title_tf.paragraphs[0].font.size = PptxPt(24)
    title_tf.paragraphs[0].font.bold = True

    body_box = slide0.shapes.add_textbox(PptxInches(0.5), PptxInches(1.0), PptxInches(9), PptxInches(2.3))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(body_tf, margin, PptxPt(2))
    paragraphs_text = _parity_unit0_paragraphs(manifest)
    for index, text in enumerate(paragraphs_text):
        paragraph = body_tf.paragraphs[0] if index == 0 else body_tf.add_paragraph()
        paragraph.text = text
        paragraph.font.size = PptxPt(11)
        paragraph.space_before = PptxPt(0)
        paragraph.space_after = PptxPt(4)
        paragraph.line_spacing = 1.0

    heading2_box = slide0.shapes.add_textbox(PptxInches(0.5), PptxInches(3.5), PptxInches(9), PptxInches(0.45))
    heading2_tf = heading2_box.text_frame
    heading2_tf.text = suite.headings[1].text
    heading2_tf.paragraphs[0].font.size = PptxPt(16)
    heading2_tf.paragraphs[0].font.bold = True

    table_fact = suite.tables[0]
    graphic_frame = slide0.shapes.add_table(
        table_fact.n_rows, table_fact.n_cols, PptxInches(0.5), PptxInches(4.15), PptxInches(6), PptxInches(2)
    )
    pptx_table = graphic_frame.table
    for cell in table_fact.cells:
        target = pptx_table.cell(cell.row, cell.col)
        target.text = cell.text
        if cell.is_header:
            for paragraph in target.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

    # Slide 1 (unit 1)
    slide1 = prs.slides.add_slide(blank_layout)
    heading3_box = slide1.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.6))
    heading3_tf = heading3_box.text_frame
    heading3_tf.text = suite.headings[2].text
    heading3_tf.paragraphs[0].font.size = PptxPt(20)
    heading3_tf.paragraphs[0].font.bold = True

    image_stream = io.BytesIO(image_bytes)
    pic_width = PptxInches(7)
    pic_height = Emu(int(pic_width * (160 / 700)))
    left = Emu(int((prs.slide_width - pic_width) / 2))
    slide1.shapes.add_picture(image_stream, left, PptxInches(1.2), width=pic_width, height=pic_height)

    caption_box = slide1.shapes.add_textbox(PptxInches(0.5), PptxInches(5.2), PptxInches(9), PptxInches(0.5))
    caption_tf = caption_box.text_frame
    caption_tf.text = suite.captions[0].text

    save_pptx(prs, out_path)


# --- stress suite --------------------------------------------------------


_LIST_STYLE_BY_INDENT = {0: "List Bullet", 1: "List Bullet 2", 2: "List Bullet 3"}


def generate_stress_docx_nested(manifest: ReferenceManifest, out_path: Path) -> None:
    fixture = manifest.stress_suite.docx_nested_structure
    doc = Document()
    for heading in fixture.headings:
        doc.add_heading(heading.text, level=heading.level)
    for item in fixture.list_items:
        style = _LIST_STYLE_BY_INDENT[item.indent_level]
        doc.add_paragraph(item.text, style=style)
    save_docx(doc, out_path)


def generate_stress_pdf_complex_layout(manifest: ReferenceManifest, out_path: Path) -> None:
    fixture = manifest.stress_suite.pdf_complex_layout
    styles = getSampleStyleSheet()
    body = styles["BodyText"]

    col1_text = next(p.text for p in fixture.paragraphs if p.column == 1)
    col2_text = next(p.text for p in fixture.paragraphs if p.column == 2)
    two_col = Table(
        [[RLParagraph(col1_text, body), RLParagraph(col2_text, body)]],
        colWidths=[3.2 * inch, 3.2 * inch],
    )
    two_col.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 0),
        ("RIGHTPADDING", (1, 0), (1, 0), 0),
    ]))

    table_fact = fixture.tables[0]
    grid: list[list[str]] = [["" for _ in range(table_fact.n_cols)] for _ in range(table_fact.n_rows)]
    header_cells: set[tuple[int, int]] = set()
    span_commands = []
    for cell in table_fact.cells:
        grid[cell.row][cell.col] = cell.text
        if cell.is_header:
            header_cells.add((cell.row, cell.col))
        if cell.col_span > 1:
            span_commands.append(("SPAN", (cell.col, cell.row), (cell.col + cell.col_span - 1, cell.row)))
    merged_table = Table(grid, colWidths=[2.1 * inch] * table_fact.n_cols)
    style_commands = [
        ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        *span_commands,
    ]
    for row, col in header_cells:
        style_commands.append(("FONTNAME", (col, row), (col, row), "Helvetica-Bold"))
    merged_table.setStyle(TableStyle(style_commands))

    story = [two_col, Spacer(1, 18), merged_table]
    doc = SimpleDocTemplate(str(out_path), pagesize=LETTER, pageCompression=0, title="", author="", subject="")
    doc.build(story, canvasmaker=_invariant_canvasmaker)


def generate_stress_pptx_overlapping_textboxes(manifest: ReferenceManifest, out_path: Path) -> None:
    fixture = manifest.stress_suite.pptx_overlapping_textboxes
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Partial physical overlap (NOT identical coordinates) so the front/back
    # stacking has a real geometric relationship, not just a z_order label.
    # Keyed by z_order: 0 = stale/behind, 1 = current/in front.
    box_position_by_z_order = {
        0: (PptxInches(1.0), PptxInches(1.0)),
        1: (PptxInches(1.8), PptxInches(1.5)),
    }
    box_size = (PptxInches(6), PptxInches(1.2))

    shapes_by_z_order = {}
    # Add lower z_order first so it renders behind; higher z_order last so
    # it renders in front -- python-pptx stacks shapes in add order.
    for box in sorted(fixture.text_boxes, key=lambda b: b.z_order):
        left, top = box_position_by_z_order[box.z_order]
        shape = slide.shapes.add_textbox(left, top, *box_size)
        shape.text_frame.text = box.text
        shapes_by_z_order[box.z_order] = shape

    # The foreground (highest z_order) box gets an opaque fill so it visibly
    # occludes the stale box behind it -- otherwise both are transparent and
    # the declared stacking has no visible effect.
    foreground_shape = shapes_by_z_order[max(shapes_by_z_order)]
    foreground_shape.fill.solid()
    foreground_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    foreground_shape.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    foreground_shape.line.width = PptxPt(1)

    table_fact = fixture.table
    graphic_frame = slide.shapes.add_table(
        table_fact.n_rows, table_fact.n_cols, PptxInches(1.0), PptxInches(3.2), PptxInches(5), PptxInches(1.5)
    )
    pptx_table = graphic_frame.table
    for cell in table_fact.cells:
        target = pptx_table.cell(cell.row, cell.col)
        target.text = cell.text
        if cell.is_header:
            for paragraph in target.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

    save_pptx(prs, out_path)


def generate_stress_pptx_native_diagram(manifest: ReferenceManifest, out_path: Path) -> None:
    fixture = manifest.stress_suite.pptx_native_diagram
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box_width = PptxInches(2)
    box_height = PptxInches(1)
    gap = PptxInches(1)
    top = PptxInches(3)

    shapes_by_id: dict[str, "pptx.shapes.autoshape.Shape"] = {}
    for index, node in enumerate(fixture.diagram_nodes):
        left = PptxInches(0.5) + index * (box_width + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, box_height)
        shape.text_frame.text = node.label
        shapes_by_id[node.fact_id] = shape

    for edge in fixture.diagram_edges:
        source_shape = shapes_by_id[edge.source]
        target_shape = shapes_by_id[edge.target]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            source_shape.left + source_shape.width, source_shape.top + source_shape.height // 2,
            target_shape.left, target_shape.top + target_shape.height // 2,
        )
        connector.begin_connect(source_shape, 3)  # right-middle connection site
        connector.end_connect(target_shape, 1)  # left-middle connection site
        connector.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
        connector.line.width = PptxPt(2)
        # Directed edge -> a visible target-end arrowhead. python-pptx has no
        # public API for this; see add_target_arrowhead()'s docstring.
        add_target_arrowhead(connector, "triangle")

    save_pptx(prs, out_path)


def generate_stress_chart_pdf(manifest: ReferenceManifest, chart_png_bytes: bytes, out_path: Path) -> None:
    canvas = _invariant_canvasmaker(str(out_path), pagesize=LETTER, pageCompression=0)
    page_width, page_height = LETTER
    image_reader_width, image_reader_height = 500, 320
    display_width = 5.0 * inch
    display_height = display_width * (image_reader_height / image_reader_width)
    x = (page_width - display_width) / 2
    y = page_height - display_height - 1.5 * inch
    canvas.drawImage(
        ImageReader(io.BytesIO(chart_png_bytes)), x, y, width=display_width, height=display_height,
        preserveAspectRatio=True, mask="auto",
    )
    canvas.showPage()
    canvas.save()


def generate_stress_scanned_pdf(manifest: ReferenceManifest, scanned_png_bytes: bytes, out_path: Path) -> None:
    """No drawString()/text calls anywhere -- the only content on the page is
    an embedded raster image, so the PDF has no digital text layer."""
    canvas = _invariant_canvasmaker(str(out_path), pagesize=LETTER, pageCompression=0)
    page_width, page_height = LETTER
    image_reader_width, image_reader_height = 900, 120
    display_width = 6.0 * inch
    display_height = display_width * (image_reader_height / image_reader_width)
    x = (page_width - display_width) / 2
    y = page_height - display_height - 1.5 * inch
    canvas.drawImage(
        ImageReader(io.BytesIO(scanned_png_bytes)), x, y, width=display_width, height=display_height,
        preserveAspectRatio=True, mask="auto",
    )
    canvas.showPage()
    canvas.save()


# --- orchestration -----------------------------------------------------


def generate_all() -> dict:
    PARITY_DIR.mkdir(parents=True, exist_ok=True)
    STRESS_DIR.mkdir(parents=True, exist_ok=True)
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)

    raw_manifest = load_manifest_raw(MANIFEST_PATH)
    manifest_sha256 = compute_manifest_sha256(raw_manifest)
    manifest = load_manifest(MANIFEST_PATH)

    console.print(f"[bold]Manifest[/bold] version={manifest.manifest_version} status={manifest.status}")
    console.print(f"[bold]manifest_sha256[/bold] (computed externally, not embedded) = {manifest_sha256}")

    diagram_png = generate_diagram_png(manifest)
    chart_png = generate_chart_png(manifest)
    scanned_png = generate_scanned_text_png(manifest)
    (IMAGES_DIR / "diagram_v1.png").write_bytes(diagram_png)
    (IMAGES_DIR / "chart_v1.png").write_bytes(chart_png)
    (IMAGES_DIR / "scanned_text_v1.png").write_bytes(scanned_png)

    generate_parity_docx(manifest, diagram_png, PARITY_DIR / "PARITY_001.docx")
    generate_parity_pdf(manifest, diagram_png, PARITY_DIR / "PARITY_001.pdf")
    generate_parity_pptx(manifest, diagram_png, PARITY_DIR / "PARITY_001.pptx")

    generate_stress_docx_nested(manifest, STRESS_DIR / "STRESS_DOCX_001.docx")
    generate_stress_pdf_complex_layout(manifest, STRESS_DIR / "STRESS_PDF_001.pdf")
    generate_stress_pptx_overlapping_textboxes(manifest, STRESS_DIR / "STRESS_PPTX_001.pptx")
    generate_stress_pptx_native_diagram(manifest, STRESS_DIR / "STRESS_PPTX_002.pptx")
    generate_stress_chart_pdf(manifest, chart_png, STRESS_DIR / "STRESS_CHART_001.pdf")
    generate_stress_scanned_pdf(manifest, scanned_png, STRESS_DIR / "STRESS_SCANNED_001.pdf")

    import hashlib

    # Exclude the report file from its own inventory -- without this, a
    # second generate_all() call (e.g. back-to-back in a determinism test,
    # with no directory cleanup in between) would pick up the PREVIOUS run's
    # report.json in the rglob scan and hash it into the NEW report, making
    # the report self-referential and stale.
    generated_files = sorted(
        path for path in GENERATED_DIR.rglob("*")
        if path.is_file() and path.name != REPORT_FILENAME
    )
    file_hashes = {}
    for path in generated_files:
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
        # POSIX-style keys (never backslashes), so the report is portable
        # and independently reproducible regardless of host OS.
        relative_key = path.relative_to(GENERATED_DIR).as_posix()
        file_hashes[relative_key] = {"sha256": digest, "size_bytes": path.stat().st_size}

    report = {
        "manifest_version": manifest.manifest_version,
        "manifest_sha256": manifest_sha256,
        "files": file_hashes,
    }
    (GENERATED_DIR / REPORT_FILENAME).write_text(json.dumps(report, indent=2, sort_keys=True), encoding="utf-8")

    console.print(f"[bold green]Generated {len(file_hashes)} files under {GENERATED_DIR}[/bold green]")
    return report


if __name__ == "__main__":
    generate_all()
